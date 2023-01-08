# Import the necessary libraries
import jira
import pptx
from datetime import date
from pptx.util import Pt

# Set up connection to Jira
jira_client = jira.JIRA(basic_auth=('your@user.name', 'YOUR_API_KEY'), server='https://your.server.address/')

# Set the names of the assignees
assignee_names = ['"John Doe"', '"Jane Doe"']

# Create an empty PowerPoint presentation
presentation = pptx.Presentation()

# Add an opening slide with a title "Sprint <sprint_number> Review"
opening_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
title = opening_slide.shapes.title
title.text = 'Daily {}'.format(date.today())

# Iterate over the assignees
for assignee in assignee_names:
  # Create a new slide for the assignee
  slide = presentation.slides.add_slide(presentation.slide_layouts[1])

  # Set the title of the slide to be the assignee's first name  (covers cases where name seperated by coma or dot)
  title = slide.shapes.title
  assignee_print = assignee.replace('"', '')
  assignee_print = assignee_print.replace(' ', '.')
  title.text = assignee_print.split('.')[0]

  shape = slide.shapes.placeholders[1]
  tf = shape.text_frame
  p = tf.add_paragraph()
  p.font.size = Pt(22)
  p.text = "In Progress:"
  p.level = 0

  # Query Jira for all tickets assigned to this assignee in the current sprint
  tickets = jira_client.search_issues('sprint in openSprints() and assignee = {} and status in ("In Progress") order by priority'.format(assignee))

  # Add a bullet point for each ticket on the slide
  bullet_slide_layout = slide.slide_layout.slide_master.slide_layouts[1]

  for ticket in tickets:
    shape = slide.shapes.placeholders[1]
    tf = shape.text_frame
    p = tf.add_paragraph()
    p.font.size = Pt(18)
    p.text = ticket.key + ': ' + ticket.fields.summary
    p.level = 1

  shape = slide.shapes.placeholders[1]
  tf = shape.text_frame
  p = tf.add_paragraph()
  p.font.size = Pt(22)
  p.text = "In review:"
  p.level = 0

  # Query Jira for all tickets assigned to this assignee in the current sprint
  tickets = jira_client.search_issues('sprint in openSprints() and assignee = {} and status in ("Review") order by priority'.format(assignee))

  # Add a bullet point for each ticket on the slide
  bullet_slide_layout = slide.slide_layout.slide_master.slide_layouts[1]

  for ticket in tickets:
    shape = slide.shapes.placeholders[1]
    tf = shape.text_frame
    p = tf.add_paragraph()
    p.font.size = Pt(18)
    p.text = ticket.key + ': ' + ticket.fields.summary
    p.level = 1

  shape = slide.shapes.placeholders[1]
  tf = shape.text_frame
  p = tf.add_paragraph()
  p.font.size = Pt(22)
  p.text = "To Do:"
  p.level = 0

  # Query Jira for all tickets assigned to this assignee in the current sprint
  tickets = jira_client.search_issues('sprint in openSprints() and assignee = {} and status in (Open, "Open (Assigned To)") order by priority'.format(assignee))

  # Add a bullet point for each ticket on the slide
  bullet_slide_layout = slide.slide_layout.slide_master.slide_layouts[1]

  for ticket in tickets:
    shape = slide.shapes.placeholders[1]
    tf = shape.text_frame
    p = tf.add_paragraph()
    p.font.size = Pt(18)
    p.text = ticket.key + ': ' + ticket.fields.summary
    p.level = 1

# Save the PowerPoint presentation
presentation.save('daily {}.pptx'.format(date.today()))
