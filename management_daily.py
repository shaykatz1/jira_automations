# Import the necessary libraries
import jira
import pptx
from datetime import date
from pptx.util import Pt


# Set up connection to Jira
jira_client = jira.JIRA(basic_auth=('your@user.name', 'YOUR_API_KEY'), server='https://your.server.address/')

# Set the names of the assignees
assignee_names = ['"John Doe"', '"Jane Doe"']

# Create an empty PowerPoint presentation
presentation = pptx.Presentation()

# Add an opening slide with a title "Sprint <sprint_number> Review"
opening_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
title = opening_slide.shapes.title
title.text = 'Scrum Daily {}'.format(date.today())

# Create a new slide
slide = presentation.slides.add_slide(presentation.slide_layouts[1])

# Set the title of the slide
title = slide.shapes.title
title.text = "In Progress"

# Iterate over the assignees
for assignee in assignee_names:

  # Query Jira for all tickets assigned to this assignee in the current sprint
  tickets = jira_client.search_issues('sprint in openSprints() and assignee = {} and status in ("In Progress") order by priority'.format(assignee))

  # Add a bullet point for each ticket on the slide
  bullet_slide_layout = slide.slide_layout.slide_master.slide_layouts[1]

  for ticket in tickets:
    shape = slide.shapes.placeholders[1]
    tf = shape.text_frame
    p = tf.add_paragraph()
    p.font.size = Pt(18)
    p.text = ticket.key + ': ' + ticket.fields.summary
    p.level = 1

# Create a new slide
slide = presentation.slides.add_slide(presentation.slide_layouts[1])

# Set the title of the slide
title = slide.shapes.title
title.text = "Done"

for assignee in assignee_names:

  # Query Jira for all tickets assigned to this assignee in the current sprint
  tickets = jira_client.search_issues('sprint in openSprints() and assignee = {} and status in ("Review") order by priority'.format(assignee))

  # Add a bullet point for each ticket on the slide
  bullet_slide_layout = slide.slide_layout.slide_master.slide_layouts[1]

  for ticket in tickets:
    shape = slide.shapes.placeholders[1]
    tf = shape.text_frame
    p = tf.add_paragraph()
    p.font.size = Pt(18)
    p.text = ticket.key + ': ' + ticket.fields.summary
    p.level = 1

# Create a new slide
slide = presentation.slides.add_slide(presentation.slide_layouts[1])

# Set the title of the slide
title = slide.shapes.title
title.text = "To Do"

for assignee in assignee_names:

  # Query Jira for all tickets assigned to this assignee in the current sprint
  tickets = jira_client.search_issues('sprint in openSprints() and assignee = {} and status in (Open, "Open (Assigned To)") order by priority'.format(assignee))

  # Add a bullet point for each ticket on the slide
  bullet_slide_layout = slide.slide_layout.slide_master.slide_layouts[1]

  for ticket in tickets:
    shape = slide.shapes.placeholders[1]
    tf = shape.text_frame
    p = tf.add_paragraph()
    p.font.size = Pt(18)
    p.text = ticket.key + ': ' + ticket.fields.summary
    p.level = 1

# Save the PowerPoint presentation
presentation.save('Scrun daily {}.pptx'.format(date.today()))
